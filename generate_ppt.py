from typing import Dict
from pptx import Presentation
from wiki import get_data
from wikipediaapi import Wikipedia
import sys


def generate_ppt(title: str, summary: str, url: str) -> None:
    presentation = Presentation("./templates/WikiSmmary.pptx")

    layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(layout)

    title_element, content_element, url_element = slide.placeholders

    title_element.text = title
    content_element.text = summary
    url_element.text = url

    # remove the first empty slide
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[0])

    presentation.save(f"./ppt/{title.replace(' ', '_')}.pptx")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("No input")
        quit()

    wiki = Wikipedia("en")

    data = [get_data(wiki, title) for title in sys.argv[1:]]
    print

    for params in data:
        if not params:
            continue
        generate_ppt(**params)
